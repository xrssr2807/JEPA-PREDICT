"""
JEPA ECG-PPG 跨通道预测模型 — PPT 生成脚本
风格：清爽专业风 | 11 页 | 中文 | 大字体版
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# ── Color Palette ──────────────────────────────────────────────
PRIMARY   = RGBColor(0x25, 0x63, 0xEB)
SECONDARY = RGBColor(0x0F, 0x76, 0x6E)
ACCENT    = RGBColor(0xF5, 0x9E, 0x0B)
DARK      = RGBColor(0x1E, 0x29, 0x3B)
BODY      = RGBColor(0x47, 0x55, 0x69)
LIGHT_BG  = RGBColor(0xF6, 0xF9, 0xFF)
CARD_BG   = RGBColor(0xF8, 0xFA, 0xFC)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
AMBER_BG  = RGBColor(0xFE, 0xF3, 0xC7)
TEAL_BG   = RGBColor(0xF0, 0xFD, 0xFB)
BLUE_BG2  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_BG3  = RGBColor(0xEF, 0xF6, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ── Helper Functions ───────────────────────────────────────────

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_round_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_text(slide, left, top, width, height, text, size=Pt(16), color=BODY,
             bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei', spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(2)
    try: p.line_spacing = spacing
    except: pass
    return tb

def add_bullets(slide, left, top, width, height, lines, size=Pt(14), color=DARK,
                font='Microsoft YaHei', spacing=1.25, sa=Pt(4)):
    """lines: list of str. Each line gets a bullet prefix if not empty."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.strip():
            p.text = "• " + line
        else:
            p.text = ""
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = sa
        try: p.line_spacing = spacing
        except: pass
    return tb

def add_section_label(slide, left, top, text, color=PRIMARY):
    """Standard '▎xxx' section label."""
    add_text(slide, left, top, Inches(5), Inches(0.4), "▎" + text,
             size=Pt(18), color=color, bold=True)

def add_title(slide, text, subtitle=None):
    """Standard slide title with accent line."""
    add_text(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.7), text,
             size=Pt(36), color=PRIMARY, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.1), Inches(1.8), Pt(4), fill_color=PRIMARY)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.35), subtitle,
                 size=Pt(16), color=BODY)

def add_card(slide, left, top, width, height, title, bullets, title_color=PRIMARY,
             accent_color=None, body_size=Pt(13)):
    """Content card with colored left accent, title, and bullets."""
    add_round_rect(slide, left, top, width, height, fill_color=CARD_BG,
                   border_color=BORDER, border_width=Pt(0.5))
    if accent_color:
        add_rect(slide, left, top, Pt(5), height, fill_color=accent_color)
    add_text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.5), Inches(0.4),
             title, size=Pt(16), color=title_color, bold=True)
    add_bullets(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.5),
                height - Inches(0.65), bullets, size=body_size)

def add_diagram_card(slide, left, top, width, height, text, bg=WHITE, border=BORDER, tc=DARK, size=Pt(13)):
    """Small labeled block for flow diagrams."""
    add_round_rect(slide, left, top, width, height, fill_color=bg, border_color=border, border_width=Pt(0.5))
    add_text(slide, left + Inches(0.08), top + Inches(0.06), width - Inches(0.16), height - Inches(0.12),
             text, size=size, color=tc, align=PP_ALIGN.CENTER, spacing=1.15)

def add_arrow(slide, left, top):
    add_text(slide, left, top, Inches(0.35), Inches(0.35), "▶", size=Pt(12), color=BODY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(0), Inches(5), H, fill_color=BLUE_BG2)
# Decorative lines
for y, w, c in [(Inches(2.2), Inches(3.2), PRIMARY), (Inches(3.1), Inches(2.4), RGBColor(0x3B,0x82,0xF6)), (Inches(4.0), Inches(2.8), RGBColor(0x60,0xA5,0xFA))]:
    add_rect(s, Inches(0.9), y, w, Pt(3), fill_color=c)
add_text(s, Inches(5.5), Inches(1.8), Inches(7.5), Inches(1.4),
         "JEPA ECG-PPG\n跨通道预测模型", size=Pt(44), color=PRIMARY, bold=True, spacing=1.15)
add_text(s, Inches(5.5), Inches(3.5), Inches(7.5), Inches(0.5),
         "模型架构、训练过程与创新点说明", size=Pt(20), color=BODY)
add_rect(s, Inches(5.5), Inches(4.25), Inches(2.2), Pt(4), fill_color=SECONDARY)
add_text(s, Inches(5.5), Inches(4.6), Inches(7.0), Inches(0.8),
         "基于 Joint Embedding Predictive Architecture\n的跨模态生理信号自监督预训练", size=Pt(14), color=BODY, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
# Slide 2: 背景与问题
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "背景与问题")

add_section_label(s, Inches(0.8), Inches(1.7), "背景", PRIMARY)
add_card(s, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5),
         "传统方法的局限", [
             "生理信号（ECG、PPG）蕴含丰富健康信息，但标注成本高昂",
             "生成式自监督方法（如 MAE）需重建原始信号，引入像素级噪声",
             "现有方法难以有效建模 ECG → PPG 的跨模态映射关系",
         ], title_color=PRIMARY, accent_color=PRIMARY, body_size=Pt(14))

add_card(s, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.6),
         "核心挑战", [
             "同一 ECG 模式可对应多种 PPG 波形（多模态性）",
             "如何在无标签条件下学习通用的生理信号表示？",
         ], title_color=ACCENT, accent_color=ACCENT, body_size=Pt(14))

add_section_label(s, Inches(7.0), Inches(1.7), "JEPA 思路", SECONDARY)
add_rect(s, Inches(7.0), Inches(2.2), Inches(5.8), Inches(3.5), fill_color=TEAL_BG)
add_text(s, Inches(7.3), Inches(2.4), Inches(5.2), Inches(0.7),
         "ECG → 嵌入空间 → 预测 PPG 表示", size=Pt(20), color=SECONDARY, bold=True)

flow = [("ECG\n心电图", BLUE_BG2), ("嵌入空间\n预测", RGBColor(0xCC,0xF0,0xEC)), ("PPG\n脉搏波", BLUE_BG2)]
for j, (t, bg) in enumerate(flow):
    x = Inches(7.3 + j * 2.1)
    add_diagram_card(s, x, Inches(3.3), Inches(1.7), Inches(1.1), t, bg=bg, border=bg, tc=DARK, size=Pt(13))
    if j < 2:
        add_text(s, x + Inches(1.7), Inches(3.55), Inches(0.4), Inches(0.4), "→", size=Pt(24), color=BODY, align=PP_ALIGN.CENTER)

add_bullets(s, Inches(7.3), Inches(4.7), Inches(5.2), Inches(1.6), [
    "不重建原始信号，在嵌入空间中进行预测",
    "更高效：避免像素/采样点级别重建误差",
    "更鲁棒：关注高层语义特征而非低层噪声",
    "无需成对标注：利用大量无标签数据进行预训练",
], size=Pt(14))

# ═══════════════════════════════════════════════════════════════
# Slide 3: 整体架构概览
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, LIGHT_BG)
add_title(s, "整体架构概览")

# Top: ECG path
add_text(s, Inches(0.3), Inches(1.7), Inches(1.2), Inches(0.4), "ECG 信号", size=Pt(14), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
blocks_t = [("CNN Stem\n4层卷积·16×压缩", Inches(1.6)), ("Transformer\n8层 Pre-LN·16 heads", Inches(3.8)), ("上下文编码器\n梯度更新", Inches(6.0)), ("预测器\n隐变量 z", Inches(8.2))]
for text, x in blocks_t:
    add_diagram_card(s, x, Inches(1.55), Inches(2.0), Inches(1.1), text, tc=DARK, size=Pt(12))
for x in [Inches(3.6), Inches(5.8), Inches(8.0)]:
    add_arrow(s, x, Inches(1.85))

# Bottom: PPG path
add_text(s, Inches(0.3), Inches(3.5), Inches(1.2), Inches(0.4), "PPG 信号", size=Pt(14), color=SECONDARY, bold=True, align=PP_ALIGN.CENTER)
blocks_b = [("CNN Stem\n4层卷积·16×压缩", Inches(1.6)), ("Transformer\n8层 Pre-LN·16 heads", Inches(3.8)), ("目标编码器\nEMA 更新（无梯度）", Inches(6.0))]
for text, x in blocks_b:
    add_diagram_card(s, x, Inches(3.35), Inches(2.0), Inches(1.1), text, tc=DARK, size=Pt(12))
for x in [Inches(3.6), Inches(5.8)]:
    add_arrow(s, x, Inches(3.65))

# Loss box
add_round_rect(s, Inches(8.3), Inches(3.0), Inches(4.5), Inches(1.9), fill_color=WHITE, border_color=SECONDARY, border_width=Pt(2))
add_text(s, Inches(8.5), Inches(3.1), Inches(4.1), Inches(0.4), "损失计算", size=Pt(16), color=SECONDARY, bold=True)
add_text(s, Inches(8.5), Inches(3.5), Inches(4.1), Inches(1.2),
         "L = MSE(预测嵌入, 目标嵌入)\n\n预测嵌入 ← Predictor(ContextEmb, z)\n目标嵌入 ← TargetEncoder(PPG)  [无梯度]",
         size=Pt(13), color=DARK, spacing=1.3)
add_arrow(s, Inches(8.05), Inches(3.65))

# Key notes
add_text(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.4),
         "💡 两个编码器结构完全相同，仅参数更新方式不同：上下文编码器通过梯度更新，目标编码器通过 EMA 更新",
         size=Pt(14), color=BODY)
add_rect(s, Inches(2.5), Inches(5.7), Inches(8.3), Inches(0.6), fill_color=BLUE_BG3, border_color=RGBColor(0xBF,0xDB,0xFE), border_width=Pt(0.5))
add_text(s, Inches(2.7), Inches(5.73), Inches(7.9), Inches(0.5),
         "EMA 更新：θ_target = m · θ_target + (1−m) · θ_context         动量 m：0.996 → 1.0（余弦调度）",
         size=Pt(14), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Slide 4: CNN Stem + Transformer 编码器
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "CNN Stem + Transformer 编码器")

# Left: CNN Stem
add_section_label(s, Inches(0.8), Inches(1.7), "CNN Stem — 一维卷积特征提取器", PRIMARY)
cnn = [
    ("Conv1d-1", "输入 1 通道 → 输出 128 通道，卷积核 7，步长 2"),
    ("Conv1d-2", "输入 128 → 输出 256，卷积核 5，步长 2"),
    ("Conv1d-3", "输入 256 → 输出 512，卷积核 5，步长 2"),
    ("Conv1d-4", "输入 512 → 输出 512，卷积核 3，步长 2"),
]
for i, (name, spec) in enumerate(cnn):
    y = Inches(2.25 + i * 0.75)
    add_round_rect(s, Inches(0.8), y, Inches(2.6), Inches(0.6), fill_color=WHITE, border_color=PRIMARY, border_width=Pt(1.5))
    add_text(s, Inches(1.0), y + Inches(0.02), Inches(2.2), Inches(0.28), name, size=Pt(13), color=PRIMARY, bold=True)
    add_text(s, Inches(1.0), y + Inches(0.3), Inches(2.2), Inches(0.28), spec, size=Pt(11), color=BODY)
    if i < 3:
        add_text(s, Inches(1.85), y + Inches(0.55), Inches(0.5), Inches(0.2), "▼", size=Pt(10), color=BODY, align=PP_ALIGN.CENTER)

# CNN summary
add_round_rect(s, Inches(3.8), Inches(2.6), Inches(2.5), Inches(2.0), fill_color=LIGHT_BG, border_color=BORDER, border_width=Pt(0.5))
add_text(s, Inches(4.0), Inches(2.8), Inches(2.1), Inches(0.7), "输入\n3000 采样点", size=Pt(14), color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.0), Inches(3.45), Inches(2.1), Inches(0.4), "↓  总步长 = 16×", size=Pt(13), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.0), Inches(3.85), Inches(2.1), Inches(0.7), "输出\n188 时间步\n(B, 512, 188)", size=Pt(14), color=DARK, bold=True, align=PP_ALIGN.CENTER, spacing=1.2)

# Right: Transformer
add_section_label(s, Inches(6.8), Inches(1.7), "Transformer 编码器", PRIMARY)
add_round_rect(s, Inches(6.8), Inches(2.25), Inches(5.8), Inches(2.9), fill_color=CARD_BG, border_color=BORDER, border_width=Pt(0.5))
add_text(s, Inches(7.1), Inches(2.35), Inches(5.3), Inches(0.35), "Pre-LN Transformer Block（×8 层堆叠）：", size=Pt(15), color=DARK, bold=True)

for i, line in enumerate([
    "LayerNorm → Multi-Head Self-Attention (16 heads) → Residual Connection",
    "LayerNorm → FFN (512 → 2048 → 512, GELU) → Residual Connection",
]):
    y = Inches(2.85 + i * 0.8)
    add_round_rect(s, Inches(7.1), y, Inches(5.3), Inches(0.65), fill_color=WHITE, border_color=BORDER, border_width=Pt(0.5))
    add_text(s, Inches(7.3), y + Inches(0.1), Inches(4.9), Inches(0.45), line, size=Pt(13), color=DARK)

add_bullets(s, Inches(7.1), Inches(4.6), Inches(5.3), Inches(1.8), [
    "自注意力机制捕捉长距离时序依赖（P-QRS-T 波段分析）",
    "16 个注意力头并行关注信号的不同特征维度",
    "自适应平均池化 → 512 维全局表示向量",
    "兼容不同输入长度（预训练 3000 点，下游 1000 点）",
], size=Pt(13))

# ═══════════════════════════════════════════════════════════════
# Slide 5: JEPA 预测架构（核心）
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "JEPA 预测架构", subtitle="模型的核心创新")

cw, ch, cy, gap, sx = Inches(3.8), Inches(4.5), Inches(1.85), Inches(0.25), Inches(0.7)

# Card 1
add_card(s, sx, cy, cw, ch, "上下文编码器 (Context Encoder)", [
    "处理 ECG 信号，生成上下文嵌入",
    "通过反向传播接收梯度更新",
    "结构：CNN Stem + Transformer + 池化",
    "唯一参与梯度优化的编码器",
], title_color=PRIMARY, accent_color=PRIMARY, body_size=Pt(14))
add_text(s, sx + Inches(0.3), cy + Inches(2.8), Inches(3.2), Inches(0.8),
         "ECG → Embedding\n▸ 梯度更新 ▸", size=Pt(14), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# Card 2 (highlighted)
cx2 = sx + cw + gap
add_round_rect(s, cx2, cy, cw, ch + Inches(0.2), fill_color=TEAL_BG, border_color=SECONDARY, border_width=Pt(2.5))
add_text(s, cx2 + Inches(0.2), cy + Inches(0.1), Inches(3.4), Inches(0.4), "预测器 (Predictor)", size=Pt(18), color=SECONDARY, bold=True)
add_bullets(s, cx2 + Inches(0.2), cy + Inches(0.6), Inches(3.4), Inches(2.3), [
    "512维上下文 → 投影至 256维 + BN + ReLU",
    "拼接 64维隐变量 z ~ N(0, I)",
    "3层 MLP: 256→256→256→256",
    "输出：256维预测嵌入",
], size=Pt(14))
add_round_rect(s, cx2 + Inches(0.15), cy + Inches(3.3), Inches(3.5), Inches(1.2), fill_color=AMBER_BG, border_color=ACCENT, border_width=Pt(1.5))
add_text(s, cx2 + Inches(0.35), cy + Inches(3.4), Inches(3.1), Inches(1.0),
         "⭐ 隐变量机制\n\nz ~ N(0, I)，64维高斯分布\n训练时采样 4 个 z，取最小 MSE", size=Pt(13), color=DARK, spacing=1.2)

# Card 3
cx3 = sx + 2*(cw + gap)
add_card(s, cx3, cy, cw, ch, "目标编码器 (Target Encoder)", [
    "处理 PPG 信号，生成目标嵌入",
    "EMA 方式更新参数（无梯度）",
    "结构：与上下文编码器完全相同",
    "防止表示坍塌，保证嵌入稳定性",
], title_color=RGBColor(0x64,0x74,0x8B), accent_color=RGBColor(0x94,0xA3,0xB8), body_size=Pt(14))
add_text(s, cx3 + Inches(0.3), cy + Inches(2.8), Inches(3.2), Inches(0.8),
         "PPG → Embedding\n▸ EMA 更新（无梯度）", size=Pt(14), color=BODY, bold=True, align=PP_ALIGN.CENTER)

# Loss formula
add_rect(s, Inches(2.0), Inches(6.65), Inches(9.3), Inches(0.55), fill_color=BLUE_BG3)
add_text(s, Inches(2.2), Inches(6.68), Inches(8.9), Inches(0.45),
         "L_pretrain = (1/B) · Σ min_{k=1..4} || Predictor(Enc_ctx(ECG), zᵏ) − Enc_tgt(PPG) ||²",
         size=Pt(15), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Slide 6: 多隐变量采样机制
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "多隐变量采样机制")

# Left: problem
add_section_label(s, Inches(0.8), Inches(1.65), "问题：PPG 信号的多模态性", ACCENT)
add_round_rect(s, Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.6), fill_color=AMBER_BG, border_color=ACCENT, border_width=Pt(0.5))
add_bullets(s, Inches(1.1), Inches(2.3), Inches(5.0), Inches(2.3), [
    "同一 ECG 模式可对应多种不同的 PPG 波形",
    "",
    "影响因素：",
    "血管弹性、血容量变化",
    "传感器接触压力与位置偏移",
    "运动伪影、环境光干扰",
    "",
    "→ 这是一个一对多映射问题",
], size=Pt(14))

# Right: solution
add_section_label(s, Inches(7.0), Inches(1.65), "解决方案：多隐变量采样", SECONDARY)
add_round_rect(s, Inches(7.0), Inches(2.15), Inches(5.8), Inches(5.0), fill_color=TEAL_BG, border_color=SECONDARY, border_width=Pt(1))

add_text(s, Inches(7.3), Inches(2.3), Inches(5.2), Inches(0.9),
         "64维高斯隐变量 z ~ N(0, I)\n对每个训练样本独立采样 4 次：",
         size=Pt(15), color=DARK, spacing=1.3)

for i, zi in enumerate(["z₁", "z₂", "z₃", "z₄"]):
    y = Inches(3.3 + i * 0.85)
    add_round_rect(s, Inches(7.3), y, Inches(1.1), Inches(0.6), fill_color=WHITE, border_color=SECONDARY, border_width=Pt(1.5))
    add_text(s, Inches(7.35), y + Inches(0.1), Inches(1.0), Inches(0.4), zi, size=Pt(16), color=SECONDARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.6), y + Inches(0.1), Inches(1.8), Inches(0.4), "+ ContextEmb → Predictor", size=Pt(13), color=DARK)
    add_text(s, Inches(10.5), y + Inches(0.1), Inches(2.0), Inches(0.4), f"→ 计算 MSE{i+1}", size=Pt(13), color=BODY)

add_round_rect(s, Inches(7.3), Inches(6.65), Inches(5.2), Inches(0.45), fill_color=WHITE, border_color=SECONDARY, border_width=Pt(2))
add_text(s, Inches(7.4), Inches(6.68), Inches(5.0), Inches(0.38),
         "最终损失 L = min(MSE₁, MSE₂, MSE₃, MSE₄)", size=Pt(16), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Slide 7: 训练过程 — 数据预处理
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, LIGHT_BG)
add_title(s, "训练过程 — 数据预处理")

# Process flow diagram at top
add_text(s, Inches(0.8), Inches(1.5), Inches(5.0), Inches(0.35),
         "▎预处理流程", size=Pt(18), color=PRIMARY, bold=True)

# Flow: pkl → extract → normalize → pt
flow_steps = [
    ("原始数据\n.pkl 文件", "115,439 个片段\n5ch × 3000pt\n30s @ 100Hz"),
    ("提取通道", "ECG (ch0)\nPPG (ch4)"),
    ("Z-score 归一化", "逐文件逐通道\n独立计算 μ、σ"),
    ("保存 .pt", "~3.1 GB"),
]
for j, (label, desc) in enumerate(flow_steps):
    fx = Inches(0.8 + j * 3.0)
    add_round_rect(s, fx, Inches(2.0), Inches(2.6), Inches(1.1), fill_color=WHITE, border_color=PRIMARY, border_width=Pt(1.5))
    add_text(s, fx + Inches(0.1), Inches(2.05), Inches(2.4), Inches(0.4), label, size=Pt(14), color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, fx + Inches(0.1), Inches(2.45), Inches(2.4), Inches(0.55), desc, size=Pt(11), color=BODY, align=PP_ALIGN.CENTER, spacing=1.2)
    if j < 3:
        add_text(s, fx + Inches(2.6), Inches(2.3), Inches(0.4), Inches(0.4), "▶", size=Pt(16), color=PRIMARY, align=PP_ALIGN.CENTER)

# Bottom: detailed explanation cards
add_round_rect(s, Inches(0.8), Inches(3.5), Inches(5.8), Inches(3.7), fill_color=WHITE, border_color=BORDER, border_width=Pt(0.5))
add_text(s, Inches(1.1), Inches(3.6), Inches(5.3), Inches(0.4),
         "数据来源与特点", size=Pt(16), color=PRIMARY, bold=True)
add_bullets(s, Inches(1.1), Inches(4.05), Inches(5.3), Inches(3.0), [
    "训练数据：115,439 个无标签多通道生理信号片段",
    "每片段包含 5 个通道，每通道 3000 个采样点",
    "采样率 100 Hz，即 30 秒的连续记录",
    "",
    "通道选择：",
    "ECG（通道 0）：体表心电信号，反映心脏电活动",
    "PPG（通道 4）：光电容积脉搏波，反映血容量变化",
    "选择 ECG→PPG 跨模态预测作为自监督预训练任务",
], size=Pt(13), spacing=1.25)

add_round_rect(s, Inches(7.0), Inches(3.5), Inches(5.8), Inches(3.7), fill_color=WHITE, border_color=BORDER, border_width=Pt(0.5))
add_text(s, Inches(7.3), Inches(3.6), Inches(5.3), Inches(0.4),
         "Z-score 归一化说明", size=Pt(16), color=SECONDARY, bold=True)
add_bullets(s, Inches(7.3), Inches(4.05), Inches(5.3), Inches(3.0), [
    "归一化方式：逐文件、逐通道独立进行",
    "公式：x_norm = (x − μ_ch) / σ_ch",
    "  μ_ch：该文件该通道所有采样点的均值",
    "  σ_ch：该文件该通道所有采样点的标准差",
    "",
    "核心目的：",
    "消除不同采集设备、不同个体之间的量纲差异",
    "同时保留信号内部的相对幅值变化特征",
    "（如 QRS 波群相对高度、脉搏波峰谷比等关键形态信息）",
], size=Pt(13), spacing=1.25)

# ═══════════════════════════════════════════════════════════════
# Slide 8: 训练过程 — 自监督预训练详解
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "训练过程 — 自监督预训练")

# Top: task definition + loss function
add_round_rect(s, Inches(0.5), Inches(1.45), Inches(7.5), Inches(2.8), fill_color=TEAL_BG, border_color=SECONDARY, border_width=Pt(1.5))
add_text(s, Inches(0.8), Inches(1.55), Inches(6.9), Inches(0.4),
         "预训练任务与损失函数", size=Pt(18), color=SECONDARY, bold=True)

add_text(s, Inches(0.8), Inches(2.0), Inches(6.9), Inches(0.6),
         "任务：给定 ECG 信号，预测 PPG 信号在嵌入空间中的表示",
         size=Pt(14), color=DARK, bold=True)

add_text(s, Inches(0.8), Inches(2.5), Inches(6.9), Inches(1.6),
         "损失函数（多隐变量 MSE）：\n\n"
         "  L_pretrain = (1/B) · Σᵦ min_{k=1..4} || Predictor(Enc_ctx(ECGᵦ), zᵦᵏ) − Enc_tgt(PPGᵦ) ||²₂\n\n"
         "  • Enc_ctx: 上下文编码器（处理 ECG，接收梯度）\n"
         "  • Enc_tgt: 目标编码器（处理 PPG，EMA 更新，无梯度）\n"
         "  • zᵦᵏ ~ N(0, I): 第 k 个隐变量采样（共 4 个），取损失最小的那个",
         size=Pt(11), color=DARK, spacing=1.2)

# Right: optimizer config
add_round_rect(s, Inches(8.3), Inches(1.45), Inches(4.7), Inches(2.8), fill_color=WHITE, border_color=PRIMARY, border_width=Pt(1.5))
add_text(s, Inches(8.55), Inches(1.55), Inches(4.2), Inches(0.4),
         "优化器配置", size=Pt(18), color=PRIMARY, bold=True)
add_bullets(s, Inches(8.55), Inches(2.05), Inches(4.2), Inches(2.0), [
    "优化器：AdamW",
    "学习率：2.5 × 10⁻³",
    "  （与 batch size 线性缩放）",
    "β₁ = 0.9, β₂ = 0.95",
    "权重衰减：0.05",
    "梯度裁剪：max_norm = 1.0",
    "Batch Size：310",
    "训练轮数：100 Epochs",
    "  （约 372 steps/epoch）",
], size=Pt(12), spacing=1.2)

# Bottom section: two cards side by side
# Left: LR schedule
add_round_rect(s, Inches(0.5), Inches(4.55), Inches(6.1), Inches(2.7), fill_color=WHITE, border_color=PRIMARY, border_width=Pt(1.5))
add_text(s, Inches(0.8), Inches(4.65), Inches(5.5), Inches(0.4),
         "学习率调度：Warmup + Cosine Annealing", size=Pt(16), color=PRIMARY, bold=True)

# Two sub-cards inside
add_round_rect(s, Inches(0.8), Inches(5.15), Inches(2.7), Inches(1.9), fill_color=BLUE_BG3, border_color=PRIMARY, border_width=Pt(0.5))
add_text(s, Inches(0.95), Inches(5.22), Inches(2.4), Inches(0.35),
         "Warmup 阶段（前 15 Epochs）", size=Pt(13), color=PRIMARY, bold=True)
add_bullets(s, Inches(0.95), Inches(5.6), Inches(2.4), Inches(1.3), [
    "学习率：10⁻⁶ → 2.5×10⁻³",
    "线性递增",
    "目的：避免训练初期",
    "梯度爆炸，让模型参数",
    "逐步适应优化方向",
], size=Pt(11), spacing=1.15)

add_round_rect(s, Inches(3.65), Inches(5.15), Inches(2.8), Inches(1.9), fill_color=BLUE_BG3, border_color=PRIMARY, border_width=Pt(0.5))
add_text(s, Inches(3.8), Inches(5.22), Inches(2.5), Inches(0.35),
         "Cosine Decay 阶段（后 85 Epochs）", size=Pt(13), color=PRIMARY, bold=True)
add_bullets(s, Inches(3.8), Inches(5.6), Inches(2.5), Inches(1.3), [
    "学习率：2.5×10⁻³ → ~0",
    "余弦曲线衰减",
    "目的：训练后期收敛",
    "到更平坦的极小值点，",
    "提升模型泛化能力",
], size=Pt(11), spacing=1.15)

# Right: EMA schedule
add_round_rect(s, Inches(6.85), Inches(4.55), Inches(6.15), Inches(2.7), fill_color=WHITE, border_color=SECONDARY, border_width=Pt(1.5))
add_text(s, Inches(7.1), Inches(4.65), Inches(5.7), Inches(0.4),
         "EMA 动量调度（目标编码器）", size=Pt(16), color=SECONDARY, bold=True)

add_text(s, Inches(7.1), Inches(5.1), Inches(5.7), Inches(0.7),
         "m(t) = m_end + (m_start − m_end) × 0.5 × (1 + cos(π · t / T))\n"
         "m_start = 0.996  →  m_end = 1.0（余弦调度）",
         size=Pt(12), color=DARK, spacing=1.25)

add_bullets(s, Inches(7.1), Inches(5.85), Inches(5.7), Inches(1.3), [
    "训练初期 m ≈ 0.996：目标编码器缓慢跟随上下文编码器",
    "  提供稳定的学习目标，有效防止表示坍塌",
    "训练末期 m → 1.0：目标编码器几乎不再更新参数",
    "  模型趋于收敛，预测器被迫学习更鲁棒的表示",
    "余弦曲线特性：初期变化平缓（利于稳定学习），末期加速收敛",
], size=Pt(11), spacing=1.15)

# ═══════════════════════════════════════════════════════════════
# Slide 9: 训练过程 — 下游微调
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, LIGHT_BG)
add_title(s, "训练过程 — 下游微调")

# Two columns: strategy on left, details on right
add_section_label(s, Inches(0.8), Inches(1.65), "下游任务：先天性心脏病 (CHD) 分类", RGBColor(0x63, 0x66, 0xF1))

# Left card: Linear Probe
add_round_rect(s, Inches(0.8), Inches(2.25), Inches(5.8), Inches(2.1), fill_color=WHITE, border_color=PRIMARY, border_width=Pt(2))
add_text(s, Inches(1.1), Inches(2.35), Inches(5.3), Inches(0.4),
         "① 线性探测（Linear Probe）— 10 Epochs", size=Pt(18), color=PRIMARY, bold=True)
add_bullets(s, Inches(1.1), Inches(2.85), Inches(5.3), Inches(1.3), [
    "冻结预训练编码器所有参数，不参与梯度更新",
    "仅训练一个随机初始化的分类头（全连接层）",
    "目的：快速、低成本地验证预训练表示的质量",
    "评估自监督学习是否学到了有区分度的生理信号特征",
], size=Pt(14), spacing=1.3)

# Right card: Full Fine-tune
add_round_rect(s, Inches(7.0), Inches(2.25), Inches(5.8), Inches(2.1), fill_color=WHITE, border_color=SECONDARY, border_width=Pt(2))
add_text(s, Inches(7.3), Inches(2.35), Inches(5.3), Inches(0.4),
         "② 全参数微调（Full Fine-tune）— 40 Epochs", size=Pt(18), color=SECONDARY, bold=True)
add_bullets(s, Inches(7.3), Inches(2.85), Inches(5.3), Inches(1.3), [
    "解冻编码器和分类头的所有参数",
    "以预训练权重初始化，较低学习率 1×10⁻⁴ 端到端训练",
    "目的：在小规模标注数据上精调，最大化下游性能",
    "借助预训练学到的通用表示，加速收敛、提升精度",
], size=Pt(14), spacing=1.3)

# Bottom: input modes
add_round_rect(s, Inches(0.8), Inches(4.65), Inches(12.0), Inches(2.6), fill_color=CARD_BG, border_color=BORDER, border_width=Pt(0.5))
add_text(s, Inches(1.1), Inches(4.75), Inches(11.5), Inches(0.4),
         "输入模式与通道配置", size=Pt(16), color=DARK, bold=True)

# Three mode cards
modes = [
    ("单通道（仅 ECG）", "仅使用 ECG 编码器\n提取心电特征 → 分类", PRIMARY),
    ("单通道（仅 PPG）", "仅使用 PPG 编码器\n提取脉搏波特征 → 分类", SECONDARY),
    ("双通道融合（ECG + PPG）", "两个编码器输出拼接\n→ 分类器\n实现跨模态信息互补", RGBColor(0x63, 0x66, 0xF1)),
]
for j, (mtitle, mdesc, mc) in enumerate(modes):
    mx = Inches(1.1 + j * 4.0)
    add_round_rect(s, mx, Inches(5.25), Inches(3.7), Inches(1.8), fill_color=WHITE, border_color=mc, border_width=Pt(1.5))
    add_text(s, mx + Inches(0.2), Inches(5.35), Inches(3.3), Inches(0.4), mtitle, size=Pt(14), color=mc, bold=True)
    add_text(s, mx + Inches(0.2), Inches(5.8), Inches(3.3), Inches(1.1), mdesc, size=Pt(13), color=DARK, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
# Slide 10: 五大创新点（内容补全 + 大字体）
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "五大创新点")

innovations = [
    ("1", "JEPA 跨模态\n预测应用", [
        "首次将 JEPA 引入 ECG→PPG",
        "跨通道生理信号预测任务",
        "嵌入空间预测：避免像素级",
        "重建误差和低层噪声干扰",
        "EMA 目标编码器防止表示坍塌",
    ], PRIMARY),
    ("2", "多隐变量\n采样机制", [
        "64维高斯隐变量 z ~ N(0, I)",
        "4次采样取最小 MSE 损失",
        "捕捉一对多映射关系",
        "（同一ECG→多种可能PPG）",
        "隐变量编码不可观测上下文",
    ], SECONDARY),
    ("3", "CNN+Transformer\n混合编码器", [
        "CNN Stem：4层Conv1d提取",
        "多尺度局部时频特征",
        "Transformer：8层Pre-LN",
        "捕捉长距离时序依赖",
        "自适应池化兼容不同输入长度",
    ], PRIMARY),
    ("4", "余弦 EMA\n动量调度", [
        "目标编码器动量 m：",
        "0.996 → 1.0 余弦递增",
        "训练初期快速适应变化",
        "训练末期趋于固定、加速收敛",
        "比固定动量更稳定、质量更高",
    ], SECONDARY),
    ("5", "无标签自监督\n+迁移学习", [
        "115k 无标签数据预训练",
        "（ECG→PPG自监督任务）",
        "线性探测验证表示质量",
        "全参数微调适配下游任务",
        "可扩展至EEG/EMG等信号",
    ], PRIMARY),
]

# Single-row: 5 cards, total = 5*2.35 + 4*0.18 = 12.47", fits 13.33"
card_w8 = Inches(2.35)
card_h8 = Inches(4.8)
gap8 = Inches(0.18)
row_start = (W - (5 * card_w8 + 4 * gap8)) / 2  # centered
card_y8 = Inches(1.6)

for i, (num, title, bullets, tc) in enumerate(innovations):
    x = row_start + i * (card_w8 + gap8)
    add_round_rect(s, x, card_y8, card_w8, card_h8, fill_color=WHITE, border_color=BORDER, border_width=Pt(0.5))
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.78), card_y8 + Inches(0.18), Inches(0.78), Inches(0.78))
    circle.fill.solid(); circle.fill.fore_color.rgb = tc; circle.line.fill.background()
    cp = circle.text_frame.paragraphs[0]
    cp.text = num; cp.font.size = Pt(28); cp.font.color.rgb = WHITE; cp.font.bold = True
    cp.font.name = 'Microsoft YaHei'; cp.alignment = PP_ALIGN.CENTER
    # Title
    add_text(s, x + Inches(0.1), card_y8 + Inches(1.2), card_w8 - Inches(0.2), Inches(0.9),
             title, size=Pt(16), color=tc, bold=True, align=PP_ALIGN.CENTER, spacing=1.15)
    # Bullets
    add_bullets(s, x + Inches(0.13), card_y8 + Inches(2.15), card_w8 - Inches(0.26), Inches(2.5),
                bullets, size=Pt(12))

# ═══════════════════════════════════════════════════════════════
# Slide 11: 下游微调结果
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "下游微调结果", subtitle="先天性心脏病 (CHD) 分类任务")

add_section_label(s, Inches(0.8), Inches(1.8), "两阶段微调策略", PRIMARY)

add_round_rect(s, Inches(0.8), Inches(2.3), Inches(5.3), Inches(1.7), fill_color=LIGHT_BG, border_color=PRIMARY, border_width=Pt(1.5))
add_text(s, Inches(1.1), Inches(2.4), Inches(4.8), Inches(0.35),
         "阶段一：线性探测 (Linear Probe, 10 Epochs)", size=Pt(16), color=PRIMARY, bold=True)
add_bullets(s, Inches(1.1), Inches(2.85), Inches(4.8), Inches(1.0), [
    "冻结预训练编码器，仅训练分类头",
    "快速验证预训练表示质量，轻量高效",
], size=Pt(14))

add_round_rect(s, Inches(0.8), Inches(4.25), Inches(5.3), Inches(1.7), fill_color=LIGHT_BG, border_color=SECONDARY, border_width=Pt(1.5))
add_text(s, Inches(1.1), Inches(4.35), Inches(4.8), Inches(0.35),
         "阶段二：全参数微调 (Full Fine-tune, 40 Epochs)", size=Pt(16), color=SECONDARY, bold=True)
add_bullets(s, Inches(1.1), Inches(4.8), Inches(4.8), Inches(1.0), [
    "解冻所有参数，lr=1×10⁻⁴ 端到端训练",
    "双通道融合：ECG + PPG 编码器输出拼接 → 分类器 → 跨模态互补",
], size=Pt(14))

# Right: result image
add_round_rect(s, Inches(6.5), Inches(1.8), Inches(6.3), Inches(4.3), fill_color=CARD_BG, border_color=BORDER, border_width=Pt(1))
img_path = os.path.join(r"F:\skills\jepa-ecg-ppg", "7b6c0da85e79e7d815c411c9cf02be0a.png")
if os.path.exists(img_path):
    try:
        s.shapes.add_picture(img_path, Inches(6.65), Inches(1.95), Inches(6.0), Inches(4.0))
    except:
        add_text(s, Inches(6.8), Inches(3.5), Inches(5.7), Inches(0.5),
                 "（图片插入失败，请手动插入）", size=Pt(14), color=BODY, align=PP_ALIGN.CENTER)
else:
    add_text(s, Inches(6.8), Inches(3.5), Inches(5.7), Inches(0.5),
             "（实验结果图请手动插入此区域）", size=Pt(14), color=BODY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Slide 12: 模型参数与训练资源
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_title(s, "模型参数与训练资源")

add_section_label(s, Inches(0.8), Inches(1.65), "模型参数统计", PRIMARY)
param_rows = [
    ("CNN Stem", "~2.5M", BLUE_BG2),
    ("Transformer 编码器 ×2 (context + target)", "~42M (2×21M)", RGBColor(0xBB,0xD6,0xFE)),
    ("上下文投影层 (context_proj)", "~0.13M", BLUE_BG2),
    ("目标投影层 (target_proj)", "~0.13M", BLUE_BG2),
    ("预测器 (Predictor)", "~0.33M", TEAL_BG),
]
for i, (name, val, bg) in enumerate(param_rows):
    y = Inches(2.15 + i * 0.58)
    add_rect(s, Inches(0.8), y, Inches(5.8), Inches(0.5), fill_color=bg)
    add_text(s, Inches(1.0), y + Inches(0.08), Inches(3.6), Inches(0.35), name, size=Pt(14), color=DARK)
    add_text(s, Inches(4.8), y + Inches(0.08), Inches(1.7), Inches(0.35), val, size=Pt(14), color=DARK, bold=True, align=PP_ALIGN.RIGHT)

add_rect(s, Inches(0.8), Inches(5.2), Inches(5.8), Inches(0.5), fill_color=PRIMARY)
add_text(s, Inches(1.0), Inches(5.25), Inches(2.5), Inches(0.4), "总参数量", size=Pt(16), color=WHITE, bold=True)
add_text(s, Inches(4.8), Inches(5.25), Inches(1.7), Inches(0.4), "~54.3M", size=Pt(16), color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
add_rect(s, Inches(0.8), Inches(5.78), Inches(5.8), Inches(0.42), fill_color=BLUE_BG2)
add_text(s, Inches(1.0), Inches(5.8), Inches(5.4), Inches(0.35),
         "可训练参数（接收梯度）≈ 27.0M（仅 context 编码器 + 预测器）", size=Pt(13), color=DARK)

add_section_label(s, Inches(7.2), Inches(1.65), "训练环境", SECONDARY)
env_items = [
    ("GPU", "NVIDIA GeForce RTX 4090D (24 GB)"),
    ("显存使用", "22.7 GB / 24 GB (92%)"),
    ("GPU 利用率", "100%"),
    ("训练数据量", "115,439 片段（约 960 小时）"),
    ("预计训练时长", "约 6.5 - 7 小时 (100 Epochs)"),
    ("Batch Size", "310"),
    ("训练效率", "约 37,200 样本/分钟"),
]
for i, (label, value) in enumerate(env_items):
    y = Inches(2.15 + i * 0.58)
    add_round_rect(s, Inches(7.2), y, Inches(5.5), Inches(0.5), fill_color=WHITE, border_color=BORDER, border_width=Pt(0.5))
    add_text(s, Inches(7.4), y + Inches(0.08), Inches(2.0), Inches(0.35), label, size=Pt(14), color=DARK, bold=True)
    add_text(s, Inches(9.4), y + Inches(0.08), Inches(3.1), Inches(0.35), value, size=Pt(14), color=BODY)

# ═══════════════════════════════════════════════════════════════
# Slide 13: 总结与展望
# ═══════════════════════════════════════════════════════════════
s = add_blank_slide()
add_bg(s, WHITE)
add_rect(s, Inches(0), Inches(6.2), W, Inches(1.3), fill_color=LIGHT_BG)
add_title(s, "总结与展望")

add_section_label(s, Inches(0.8), Inches(1.65), "核心价值", PRIMARY)
add_round_rect(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(3.8), fill_color=LIGHT_BG, border_color=BORDER, border_width=Pt(0.5))
add_bullets(s, Inches(1.1), Inches(2.35), Inches(5.0), Inches(3.4), [
    "利用 11.5 万无标签生理信号实现自监督预训练",
    "JEPA 嵌入空间预测，避免信号重建噪声与误差",
    "EMA 目标编码器防止表示坍塌",
    "多隐变量采样捕捉一对多映射关系",
    "CNN + Transformer 混合编码器兼顾局部与全局特征",
    "下游 CHD 分类任务验证预训练表示有效性",
], size=Pt(14))

add_section_label(s, Inches(7.0), Inches(1.65), "未来方向", SECONDARY)
add_round_rect(s, Inches(7.0), Inches(2.15), Inches(5.6), Inches(3.8), fill_color=TEAL_BG, border_color=SECONDARY, border_width=Pt(0.5))
add_bullets(s, Inches(7.3), Inches(2.35), Inches(5.0), Inches(3.4), [
    "扩展至 EEG、EMG 等其他生理信号的跨模态预训练",
    "探索更长序列的多尺度预训练策略",
    "适配可穿戴设备实时推理场景",
    "结合临床标注数据进一步验证医学价值",
    "探索多任务学习（分类 + 回归 + 生成）框架",
], size=Pt(14))

# Bottom takeaway
add_rect(s, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.65), fill_color=PRIMARY)
add_text(s, Inches(1.2), Inches(6.53), Inches(10.9), Inches(0.55),
         "JEPA 为大规模无标签生理信号的自监督学习提供了通用、高效、可扩展的解决方案",
         size=Pt(18), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output_path = r"F:\skills\jepa-ecg-ppg\JEPA_ECG-PPG_项目汇报_v6.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
